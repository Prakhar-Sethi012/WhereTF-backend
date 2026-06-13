"""
extractors.py
-------------
One extractor class / function per supported file type.

Every extractor returns a list of ``ChunkData`` dicts:

    {
        "chunk_index" : int,   # 0-based position within the file
        "content_text": str,   # full text for this chunk (page / slide / block)
        # "embedding" and "keyword_tokens" are added downstream
    }

Embedded-image OCR (PDF, DOCX, PPTX) appends OCR text to the same chunk
so that the text and its visual content share the same vector space.

Rolling/overlapping chunking is used for all text-heavy formats so that
context is never silently cut at a hard boundary.

No database sessions, no embedding calls, no side effects.
"""

from __future__ import annotations

import io
import json
import logging
import zipfile
from pathlib import Path
from typing import TypedDict

from .ocr import ocr_image_bytes

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Shared type
# ---------------------------------------------------------------------------

class ChunkData(TypedDict):
    chunk_index: int
    content_text: str
    # embedding and keyword_tokens intentionally absent here;
    # they are injected by the embedding pipeline.


# ---------------------------------------------------------------------------
# Chunking configuration
# ---------------------------------------------------------------------------

# Rolling-window parameters (characters).
# CHUNK_SIZE  : maximum characters per chunk.
# CHUNK_OVERLAP: how many characters from the end of chunk N are repeated
#                at the start of chunk N+1, preserving cross-boundary context.
CHUNK_SIZE    = 1_500
CHUNK_OVERLAP = 200


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _clean(text: str) -> str:
    """Collapse excessive whitespace without destroying newlines."""
    import re
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _rolling_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Split *text* into overlapping sliding windows.

    Each window is at most *size* characters.  The hop between windows is
    exactly ``size - overlap`` characters, so every boundary region appears
    in two successive chunks and no context is silently lost.

    The soft-break on the last newline inside each window only fires when
    it falls within the hop zone ``[start + step, end]``, so it can never
    shrink the step and cause degenerate micro-chunks.

    Returns a list of non-empty stripped strings.
    """
    if not text:
        return []

    step   = size - overlap          # guaranteed minimum forward progress
    chunks: list[str] = []
    start  = 0
    length = len(text)

    while start < length:
        end = min(start + size, length)

        # Soft-break: snap to last newline in the tail of the window,
        # but only if that keeps the step at least `step` chars long.
        if end < length:
            nl = text.rfind("\n", start + step, end)
            if nl != -1:
                end = nl + 1          # include the newline in this chunk

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start += step                 # fixed-size hop, always progresses

    return chunks


def _zip_image_bytes(zf: zipfile.ZipFile, name: str) -> bytes:
    """Read a member from an open ZipFile; return empty bytes on error."""
    try:
        return zf.read(name)
    except Exception:  # noqa: BLE001
        return b""


def _pdf_image_to_png_bytes(img_meta: dict) -> bytes:
    """
    Convert a pdfplumber image-metadata dict to a valid PNG byte-string
    that PIL can open.

    pdfplumber stores raw decoded pixel data (not a container format) in
    img_meta["stream"].get_data().  We reconstruct a proper PNG from the
    dimensions and colour-space information that pdfplumber also exposes.

    Supported colour spaces: DeviceRGB (3-channel), DeviceGray (1-channel),
    and DeviceCMYK (4-channel, converted to RGB).
    Falls back to trying the raw bytes directly (sometimes the stream IS
    a JPEG or PNG already).
    """
    try:
        from PIL import Image  # type: ignore

        stream = img_meta.get("stream")
        if stream is None:
            return b""

        raw = stream.get_data()

        # --- fast path: bytes are already a valid image container ----------
        try:
            Image.open(io.BytesIO(raw)).verify()
            return raw          # it worked — JPEG / PNG inside PDF
        except Exception:
            pass                # fall through to pixel reconstruction

        # --- reconstruct from raw pixel buffer ------------------------------
        width  = int(img_meta.get("width",  img_meta.get("Width",  0)))
        height = int(img_meta.get("height", img_meta.get("Height", 0)))
        cs     = str(img_meta.get("colorspace", img_meta.get("ColorSpace", "RGB")))

        if width == 0 or height == 0:
            return b""

        if "gray" in cs.lower() or "grey" in cs.lower():
            mode, channels = "L", 1
        elif "cmyk" in cs.lower():
            mode, channels = "CMYK", 4
        else:
            mode, channels = "RGB", 3

        expected = width * height * channels
        if len(raw) < expected:
            return b""

        pil_img = Image.frombytes(mode, (width, height), raw[:expected])

        # CMYK → RGB so EasyOCR never sees CMYK
        if mode == "CMYK":
            pil_img = pil_img.convert("RGB")

        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        return buf.getvalue()

    except Exception as exc:  # noqa: BLE001
        logger.debug("[pdf] Could not convert image to PNG: %s", exc)
        return b""


# ---------------------------------------------------------------------------
# PDF extractor
# ---------------------------------------------------------------------------

def extract_pdf(path: Path) -> list[ChunkData]:
    """
    One page = one text block.  The text blocks are then fed through the
    rolling-window chunker so long pages produce overlapping chunks.

    Each page block also includes OCR output from any embedded raster
    images on that page.
    """
    try:
        import pdfplumber  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pip install pdfplumber") from exc

    page_blocks: list[str] = []

    with pdfplumber.open(str(path)) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            parts: list[str] = []

            # 1. Native text layer
            native = page.extract_text() or ""
            if native.strip():
                parts.append(native)

            # 2. Embedded images → PNG reconstruction → OCR
            for img_meta in page.images:
                try:
                    png_bytes = _pdf_image_to_png_bytes(img_meta)
                    if not png_bytes:
                        continue
                    ocr_text = ocr_image_bytes(png_bytes)
                    if ocr_text:
                        parts.append(f"[OCR] {ocr_text}")
                except Exception as exc:  # noqa: BLE001
                    logger.debug("[pdf] OCR failed on page %d image: %s", page_idx, exc)

            block = _clean("\n".join(parts))
            if block:
                page_blocks.append(block)

    # Join all page blocks then apply rolling chunker so cross-page context
    # is preserved at page boundaries too.
    full_text = "\n\n".join(page_blocks)
    raw_chunks = _rolling_chunks(full_text)

    return [
        ChunkData(chunk_index=i, content_text=c)
        for i, c in enumerate(raw_chunks)
    ]


# ---------------------------------------------------------------------------
# DOCX extractor
# ---------------------------------------------------------------------------

def extract_docx(path: Path) -> list[ChunkData]:
    """
    All paragraphs are joined into a single document string, then split
    via the rolling-window chunker.  Embedded images are OCR-ed and
    inserted at the end of the text before chunking.
    """
    try:
        import docx as _docx  # type: ignore  (python-docx)
    except ImportError as exc:
        raise RuntimeError("pip install python-docx") from exc

    doc = _docx.Document(str(path))

    # ---- collect image bytes from word/media/ ----
    with zipfile.ZipFile(str(path), "r") as zf:
        media_bytes: list[bytes] = [
            _zip_image_bytes(zf, name)
            for name in zf.namelist()
            if name.startswith("word/media/")
        ]

    # ---- build full document text ----
    lines: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)

    # Append OCR from all embedded images
    for img_bytes in media_bytes:
        if not img_bytes:
            continue
        ocr_text = ocr_image_bytes(img_bytes)
        if ocr_text:
            lines.append(f"[OCR] {ocr_text}")

    full_text = _clean("\n".join(lines))
    raw_chunks = _rolling_chunks(full_text)

    return [
        ChunkData(chunk_index=i, content_text=c)
        for i, c in enumerate(raw_chunks)
    ]


# ---------------------------------------------------------------------------
# PPTX extractor
# ---------------------------------------------------------------------------

def extract_pptx(path: Path) -> list[ChunkData]:
    """
    One chunk per slide (slides are natural semantic units in presentations).
    Slide text + speaker notes + OCR of every embedded image on that slide.
    Long slides are further split by the rolling chunker.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pip install python-pptx") from exc

    prs = Presentation(str(path))
    chunks: list[ChunkData] = []
    global_chunk_idx = 0

    for slide in prs.slides:
        parts: list[str] = []

        # 1. Text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

        # 2. Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        # 3. Embedded images via shape relationships (most reliable path)
        seen_blobs: set[int] = set()
        for rel in slide.part.rels.values():
            try:
                if "image" not in rel.reltype:
                    continue
                blob = rel.target_part.blob
                blob_id = id(blob)
                if blob_id in seen_blobs:
                    continue
                seen_blobs.add(blob_id)
                ocr_text = ocr_image_bytes(blob)
                if ocr_text:
                    parts.append(f"[OCR] {ocr_text}")
            except Exception:  # noqa: BLE001
                pass

        slide_text = _clean("\n".join(parts))
        if not slide_text:
            continue

        for sub_chunk in _rolling_chunks(slide_text):
            chunks.append(ChunkData(chunk_index=global_chunk_idx, content_text=sub_chunk))
            global_chunk_idx += 1

    return chunks


# ---------------------------------------------------------------------------
# XLSX extractor
# ---------------------------------------------------------------------------

def extract_xlsx(path: Path) -> list[ChunkData]:
    """
    Convert each sheet to a readable text table (tab-separated).
    Each sheet becomes one rolling-chunked text block so large sheets
    produce overlapping chunks.
    """
    try:
        import openpyxl  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pip install openpyxl") from exc

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[ChunkData] = []
    global_idx = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []

        for row in ws.iter_rows(values_only=True):
            # Skip completely empty rows
            if all(cell is None for cell in row):
                continue
            cells = [str(cell) if cell is not None else "" for cell in row]
            rows.append("\t".join(cells))

        if not rows:
            continue

        # Prepend sheet name as context header
        sheet_text = _clean(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        for sub_chunk in _rolling_chunks(sheet_text):
            chunks.append(ChunkData(chunk_index=global_idx, content_text=sub_chunk))
            global_idx += 1

    wb.close()
    return chunks


# ---------------------------------------------------------------------------
# Standalone image extractor (OCR only)
# ---------------------------------------------------------------------------

def extract_image(path: Path) -> list[ChunkData]:
    """Single chunk: full OCR output of the image file."""
    img_bytes = path.read_bytes()
    ocr_text = ocr_image_bytes(img_bytes)
    if not ocr_text:
        return []
    # Images rarely exceed one chunk but run through rolling chunker for safety
    return [
        ChunkData(chunk_index=i, content_text=c)
        for i, c in enumerate(_rolling_chunks(_clean(ocr_text)))
    ]


# ---------------------------------------------------------------------------
# Plain-text / code extractor
# ---------------------------------------------------------------------------

def extract_text(path: Path) -> list[ChunkData]:
    """
    Read a UTF-8 text file and split it into rolling overlapping chunks.
    JSON files are pretty-printed before chunking.
    CSV files get a column-header prefix on every chunk for LLM context.
    """
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        logger.warning("[text] Cannot read %s: %s", path, exc)
        return []

    ext = path.suffix.lower()

    # --- JSON: pretty-print for readability ---------------------------------
    if ext == ".json":
        try:
            raw = json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            pass

    # --- CSV: prepend header to every chunk for grounding -------------------
    if ext == ".csv":
        lines = raw.splitlines()
        header = lines[0] if lines else ""
        body   = "\n".join(lines[1:]) if len(lines) > 1 else ""
        sub_chunks = _rolling_chunks(_clean(body))
        return [
            ChunkData(chunk_index=i, content_text=_clean(f"[Columns: {header}]\n{c}"))
            for i, c in enumerate(sub_chunks)
        ] if sub_chunks else []

    raw_chunks = _rolling_chunks(_clean(raw))
    return [
        ChunkData(chunk_index=i, content_text=c)
        for i, c in enumerate(raw_chunks)
    ]


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

#: Map every supported extension to its extractor function.
#: Add new extensions here; no other file needs changing.
_EXTENSION_MAP: dict[str, object] = {
    # ── Documents ─────────────────────────────────────────────────────────
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,   # openpyxl handles legacy xls via compatibility
    # ── Images (standalone OCR) ───────────────────────────────────────────
    ".png":  extract_image,
    ".jpg":  extract_image,
    ".jpeg": extract_image,
    ".webp": extract_image,
    ".bmp":  extract_image,
    ".tiff": extract_image,
    ".tif":  extract_image,
    # ── Plain text & markup ───────────────────────────────────────────────
    ".txt":  extract_text,
    ".md":   extract_text,
    ".rst":  extract_text,
    ".tex":  extract_text,
    # ── Data formats ─────────────────────────────────────────────────────
    ".json": extract_text,
    ".csv":  extract_text,
    ".tsv":  extract_text,
    ".xml":  extract_text,
    ".yaml": extract_text,
    ".yml":  extract_text,
    ".toml": extract_text,
    # ── Python ───────────────────────────────────────────────────────────
    ".py":   extract_text,
    ".pyi":  extract_text,
    ".ipynb": extract_text,
    # ── JavaScript / TypeScript ───────────────────────────────────────────
    ".js":   extract_text,
    ".jsx":  extract_text,
    ".ts":   extract_text,
    ".tsx":  extract_text,
    ".mjs":  extract_text,
    ".cjs":  extract_text,
    # ── Web ───────────────────────────────────────────────────────────────
    ".html": extract_text,
    ".htm":  extract_text,
    ".css":  extract_text,
    ".scss": extract_text,
    ".sass": extract_text,
    # ── Systems / compiled languages ─────────────────────────────────────
    ".c":    extract_text,
    ".h":    extract_text,
    ".cpp":  extract_text,
    ".cc":   extract_text,
    ".cxx":  extract_text,
    ".hpp":  extract_text,
    ".hxx":  extract_text,
    ".cs":   extract_text,   # C#
    ".java": extract_text,
    ".kt":   extract_text,   # Kotlin
    ".swift": extract_text,
    ".go":   extract_text,
    ".rs":   extract_text,   # Rust
    ".zig":  extract_text,
    # ── Scripting / shell ─────────────────────────────────────────────────
    ".sh":   extract_text,
    ".bash": extract_text,
    ".zsh":  extract_text,
    ".fish": extract_text,
    ".ps1":  extract_text,   # PowerShell
    ".bat":  extract_text,
    ".cmd":  extract_text,
    # ── Ruby / PHP / others ───────────────────────────────────────────────
    ".rb":   extract_text,
    ".php":  extract_text,
    ".lua":  extract_text,
    ".pl":   extract_text,   # Perl
    ".r":    extract_text,   # R
    ".scala": extract_text,
    ".ex":   extract_text,   # Elixir
    ".exs":  extract_text,
    ".erl":  extract_text,   # Erlang
    ".hs":   extract_text,   # Haskell
    ".ml":   extract_text,   # OCaml
    ".clj":  extract_text,   # Clojure
    # ── Config / infra ────────────────────────────────────────────────────
    ".ini":  extract_text,
    ".cfg":  extract_text,
    ".conf": extract_text,
    ".env":  extract_text,
    ".dockerfile": extract_text,
    ".tf":   extract_text,   # Terraform
    ".hcl":  extract_text,
    ".sql":  extract_text,
    ".graphql": extract_text,
    ".proto": extract_text,  # Protobuf
}


def extract(path: Path) -> list[ChunkData]:
    """
    Dispatch to the correct extractor based on file extension.

    Parameters
    ----------
    path:
        Resolved :class:`~pathlib.Path` to the file.

    Returns
    -------
    list[ChunkData]
        Zero or more chunk dicts ready for the embedding pipeline.
    """
    # Handle extensionless files often found in codebases (Makefile, Dockerfile…)
    ext = path.suffix.lower() or path.name.lower()
    extractor = _EXTENSION_MAP.get(ext)

    # Special-case common extensionless names
    if extractor is None and path.name.lower() in {
        "makefile", "dockerfile", "jenkinsfile", "vagrantfile",
        "gemfile", "rakefile", "procfile", "brewfile",
        ".gitignore", ".gitattributes", ".editorconfig",
        "requirements", "pipfile", "cargo.lock", "go.sum",
    }:
        extractor = extract_text

    if extractor is None:
        logger.debug("[extractor] No extractor for '%s' (%s)", ext, path.name)
        return []

    try:
        return extractor(path)
    except Exception as exc:  # noqa: BLE001
        logger.error("[extractor] Failed to extract %s: %s", path, exc)
        return []